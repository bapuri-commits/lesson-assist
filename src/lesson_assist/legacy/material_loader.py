"""강의자료(PDF/PPT) 텍스트 추출 및 RAG 저장.

eclass_crawler가 다운로드한 PDF/PPT 파일에서 텍스트를 추출하여
ChromaDB에 저장한다. 요약 시 "슬라이드에 있는 것처럼" 참조 가능.
"""
from __future__ import annotations

from pathlib import Path

from loguru import logger


def extract_pdf_text(pdf_path: Path) -> list[str]:
    """PDF에서 페이지별 텍스트를 추출한다."""
    import pdfplumber

    pages: list[str] = []
    with pdfplumber.open(str(pdf_path)) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text() or ""
            text = text.strip()
            if text:
                pages.append(f"[p.{i + 1}] {text}")
    return pages


def extract_pptx_text(pptx_path: Path) -> list[str]:
    """PPT에서 슬라이드별 텍스트를 추출한다."""
    from pptx import Presentation

    slides: list[str] = []
    prs = Presentation(str(pptx_path))
    for i, slide in enumerate(prs.slides):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
        if texts:
            combined = " ".join(texts)
            slides.append(f"[slide {i + 1}] {combined}")
    return slides


def extract_text(file_path: Path) -> list[str]:
    """파일 확장자에 따라 텍스트를 추출한다."""
    suffix = file_path.suffix.lower()
    try:
        if suffix == ".pdf":
            return extract_pdf_text(file_path)
        elif suffix == ".pptx":
            return extract_pptx_text(file_path)
        elif suffix == ".ppt":
            logger.warning(f".ppt 구형 포맷 미지원 (python-pptx는 .pptx만 지원): {file_path.name}")
            return []
        else:
            logger.debug(f"텍스트 추출 미지원 형식: {file_path.name}")
            return []
    except Exception as e:
        logger.warning(f"텍스트 추출 실패: {file_path.name} — {e}")
        return []


def extract_and_store_materials(
    store,
    course: str,
    material_paths: list[Path],
) -> int:
    """강의자료 파일들에서 텍스트를 추출하여 RAG에 저장한다.

    이미 저장된 파일은 건너뛴다 (파일명 기반 중복 방지).

    Returns:
        새로 저장된 총 청크 수.
    """
    total_chunks = 0

    for path in material_paths:
        if not path.is_file():
            continue

        pages = extract_text(path)
        if not pages:
            continue

        try:
            count = store.add_material(course, pages, path.name)
            total_chunks += count
        except Exception as e:
            logger.warning(f"자료 RAG 저장 실패: {path.name} — {e}")

    if total_chunks > 0:
        logger.info(f"eclass 자료 RAG 저장: {course} → {total_chunks}개 청크 ({len(material_paths)}개 파일)")
    return total_chunks
